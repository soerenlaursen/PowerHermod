#!/usr/local/bin/python2.7
# -*- coding: utf-8 -*-
'''
PowerHermod.Grane -- shortdesc

PowerHermod.Grane is a description

It defines classes_and_methods

We are using the microsoft powerpoint module for python (python-pptx):
https://python-pptx.readthedocs.org/en/latest/

@author:     Soeren Laursen

@copyright:  2014, 2015, 2016, 2017 FCOO. All rights reserved.

@license:    See LICENSE file.

@contact:    sln@fcoo.dk
@deffield    updated: Updated

@install:
On Ubuntu 14.04 LTS
sudo apt-get install python-pip
sudo apt-get install python-pil
sudo apt-get install python-lxml
sudo apt-get install python-yaml
sudo apt-get install xvfb
sudo apt-get install cutycapt
sudo pip install python-pptx

Sometimes the powerpoint can be very large, increase postfix email size:
postconf -e 'message_size_limit = 104857600'

On Ubuntu 16.04 LTS
- Not testet yet!

Notes:
Basic concept is that a briefing (be it pptx, odp, or pdf), is a collection of images, wrapped in a container.
Source can be an image path or an url, which is captured and saved as an image, then parsed into the briefing.

If image is created from url, it will be requested at the proper resolution, but if one wants it scaled (smaller/larger than "reality") the imageScale parameter can be used.

Image layout is based on a 6x4 grid (so from 1 to 24 images per slide).
An image location and size can be defined by an identifier, which is a combination of the upper left corner and the desired size.
Some examples:
1,1 4x3 : full slide
1,1 2x3 : left half
3,1 2x3 : right half
4,1 1x1 : top right corner (1x1)
3,2 2x2 : top right corner (2x2)

Since a legend is usually needed (and often shared between iamges it can be located on either the left (l), bottom (b), or right (r) side of the grid:
1x3l : left side
2x1r : middle right
2x2b : middle bottom

Images will be generated in the correct aspect ratio if possible otherwise scaled to fit (preserving the aspect ratio).
'''

# The path to cutycapt program we use to generate images from a WMS service
cutycaptPath = '/usr/bin/cutycapt'

#
from pptx import Presentation
from pptx.util import Inches, Px, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import sys
import os
# The glob library
import glob

# To fetch an url
import urllib

# Imports for sending email
import subprocess

# Imports for sending email
import smtplib

# Imports for base64 encode/decode
import base64

#
import unicodedata

# Imports for creating MIME encoded email
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.MIMEBase import MIMEBase
from email.MIMEText import MIMEText
from email.Utils import COMMASPACE, formatdate
from email import Encoders

# Import for logging
import logging

# Split parameters
import shlex

# Import for date calculation/handling
#from datetime import datetime, time, date, timedelta
import datetime

#from time import strftime

from argparse import ArgumentParser
from argparse import RawDescriptionHelpFormatter

# xvfb python bindings
from xvfbwrapper import Xvfb

#  yaml support
import yaml

# zipfile support
# We are planing to use a zipfile as a container for png files.

__all__ = []
__version__ = 0.1
__date__ = '2014-04-08'
__updated__ = '2016-03-10'

DEBUG = 0
TESTRUN = 0
PROFILE = 0


def fetchImage(url,filename):
    urllib.urlretrieve(url, filename)

def generateImageFromWMS(url, filename, xdisplay, screenSize='medium', username = '', password = '', imageResolution=None, delay=5000):
    # Examples of resolutions
    # 16:9
    # 1024×576

    # 4:3
    # 800x600
    # 1600x1200

    resolution = {}
    if imageResolution:
        resolution['input'] = imageResolution
        screenSize = 'input' # overrides screenSize
    else:
        resolution['small']  = (800,600)   # legacy screen (very poor resolution)
        resolution['medium'] = (1024,768)  # normal (old) projector or screen
        resolution['large']  = (1500,1100) # modern projector or screen (full hd)
        resolution['huge']   = (4000,4000) # high-resolution projektor or screen (4k)

    logger.debug('Generate image from : ' + url + ' to file : ' + filename )

    # PHP Code
    #http_auth  = "Authorization: Basic " . base64_encode("{$user_id}:{$password}");

    # Do we have a user/password?
    if len(username) != 0:
        http_auth_string  = "Authorization: Basic " + base64.b64encode(username + ':' + password)
        cmd = '''{cutycapt} -display :{x} --header="{http_auth}" --url='{u}' --min-width={w} --min-height={h} --delay={d} --javascript=on --out={f}'''.format(cutycapt = cutycaptPath, u = url, x = xdisplay, d = delay, f = filename, w = resolution[screenSize][0], h = resolution[screenSize][1], http_auth = http_auth_string )
    else:
        cmd = '''{cutycapt} -display :{x}  --url='{u}' --min-width={w} --min-height={h} --delay={d} --javascript=on --out={f}'''.format(cutycapt = cutycaptPath, u = url, x = xdisplay, d = delay, f = filename, w = resolution[screenSize][0], h = resolution[screenSize][1] )

    logger.debug('Command : ' + cmd)
    print("COMMAND : " + cmd)

    proc = subprocess.Popen(shlex.split(cmd))
    proc.communicate()

def formatFilename( filename, epoch ):
    '''
    Formats a filename (with strftime) from the current epoch.
    :param filename:
    :param epoch:
    '''
    ufilename = unicode( filename )
    filename = unicodedata.normalize('NFKD', ufilename).encode('ascii', 'ignore')
    return epoch.strftime(filename)

def formatUrl( url,epoch ):
    '''
    Format an url from the giving epoch, using strftime.
    :param url:
    :param epoch:
    '''

#SHOULD WORK BUT WMS ISN'Y USING UTF8
    return epoch.strftime( url.encode('utf-8') )
    # Danish Kludge
#    url = url.replace('','æ')
#    return epoch.strftime( url )

def formatTitle( title, epoch ):
    '''
    Format a title from the giving epoch, using strftime.
    :param title:
    :param epoch:
    '''
    title = title.replace('_',' ').encode('utf-8')
    return epoch.strftime( title )

def formatPageTitle( title, epoch ):
    '''
    Format a page title from the giving epoch, using strftime.
    Replace _ with ' '
    :param title:
    :param epoch:
    '''
    return epoch.strftime(title.replace('_',' ').encode('utf-8'))

def get_storage_folder(destdir, name, t):
    """Creates sub directories when needed."""
    year = t.strftime('%Y')
    month = t.strftime('%m')
    day = t.strftime('%d')
    cdir = os.path.join(os.path.normpath(destdir), year, month, day)
    try:
        os.makedirs(cdir)
    except:
        # Throws an exception if dir already exists
        pass
    logging.info('file: ' + str(os.path.join(cdir, name) ) )
    pathAndFilename = os.path.join(cdir, name)
    return pathAndFilename

def downloadImage(url, filename, username, password ):
    return "lll"

def readAndVerifyEmailRecipient( recipientStringList ):
    recipientSplit = recipientStringList.split(",")
    recipientArray = []
    for recipient in recipientSplit:
        if (len(recipient.strip()) > 3) & (recipient.strip().find("@") != -1) :
            recipientArray.append( recipient.strip())
    return recipientArray


def processHermodSubscription( hermonConfiguration ):
    startdate = hermodConfiguration['startdate']
    enddate = hermodConfiguration['enddate']

    # Convert the start- and enddate to a datetime periode
    startdate = datetime.datetime.combine(startdate, datetime.datetime.min.time())
    enddate = datetime.datetime.combine(enddate, datetime.datetime.max.time())

    # We are using UTC as internal time
    utcNow = datetime.datetime.utcnow()
    currentYear = utcNow.year
    currentMonth = utcNow.month
    currentDay = utcNow.day
    currentHour = utcNow.hour

    # Are we suppose to delivere now?
    if startdate <= utcNow <= enddate:
        logger.info("We are delivering data")

    # Convert the selected epoch to a list of integer
    selectedEpoch = hermodConfiguration['epoch']
    #epochList = [int(x) for x in selectedEpoch.split(",")]
    epochList = [int(selectedEpoch)]

    if 'startEpoch' in hermodConfiguration.keys():
        startEpoch = utcNow + datetime.timedelta(hours=(hermodConfiguration['startEpoch']-currentHour))
    else:
        startEpoch = utcNow

    # Debug - use this to simulate a epoch hit.
    #currentHour = 10

    # Do we have a match on epoch or are we prescribing time
    if currentHour in epochList:
        pages = hermodConfiguration['pages']

        # Read standard parameters
        if 'timestampPattern' in hermodConfiguration.keys():
            filetimestampPattern = hermodConfiguration['timestampPattern']
        else:
            filetimestampPattern = '-%Y-%m-%dT%H0000'

        if 'timestampPattern' in hermodConfiguration.keys():
            headlinePattern = hermodConfiguration['headlinePattern']
        else:            headlinePattern = ' - %Y-%m-%dT%H0000'

        if 'template' in hermodConfiguration.keys():
            powerpointTemplate = hermodConfiguration['template']
            powerpointTemplate = os.path.join(os.path.normpath(powerhermodConfigurationFiles), powerpointTemplate)
        else:
            logger.error('Template not defined!')
            exit(-1)

        if 'filename' in hermodConfiguration.keys():
            outputFilename = hermodConfiguration['filename']
        else:
            logger.error('Output filename not defined!')
            exit(-1)

        pptTitle = "Not defined!"
        if 'description' in hermodConfiguration.keys():
            pptTitle = hermodConfiguration['description']

        emailtext = ''
        if 'emailtext' in hermodConfiguration.keys():
            emailtext = hermodConfiguration['emailtext']

        # Where do we keep scratch file and files for the archive
        # Default using /tmp
        PowerHermodScratch = '/tmp'
        PowerHermodArchive = '/tmp'

        if 'scratchPath' in hermodConfiguration.keys():
            PowerHermodScratchPath = hermodConfiguration['scratchPath']

        if 'archivePath' in hermodConfiguration.keys():
            PowerHermodArchivePath = hermodConfiguration['archivePath']

        if 'attachImages' in hermodConfiguration.keys():
            attachImages = hermodConfiguration['attachImages']
        else:
            attachImages = False

        if 'compressImages' in hermodConfiguration.keys():
            compressImages = hermodConfiguration['compressImages']
        else:
            compressImages = True

        if 'delay' in hermodConfiguration.keys():
            globalDelay = hermodConfiguration['delay'] * 1000 # convert to microseconds
        else:
            globalDelay = 5 * 1000

        prs = Presentation( powerpointTemplate )

        # Start a new X framebuffer
        vdisplay = Xvfb()
        vdisplay.start()

        # Title page
        blank_slidelayout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slidelayout)
        title = slide.shapes.title
        title.text = pptTitle
        subtitle = slide.shapes.placeholders[1]
        subtitle.text = utcNow.strftime('METOC - %Y-%m-%dT%H:00:00 UTC' )

        pages = hermodConfiguration['pages']
        # Picture
        #    pages = []
        # This is in UTC
        globalStartEpoch = utcNow
#        globalStartEpoch = currentHour

        if 'pageLayout' in hermodConfiguration.keys():
            pageLayout = map(int,hermodConfiguration['pageLayout'].split('x'))
        else:
            pageLayout = 6, 4

#            pageWidth  = 27.5 # could possible computed automatically from template???
#            pageHeight = 19.0 # could possible computed automatically from template???
            pageWidth  = 25.5 # could possible computed automatically from template???
            pageHeight = 19.0 # could possible computed automatically from template???

        if 'screenSize' in hermodConfiguration.keys():
            screenSize = hermodConfiguration['screenSize']
        else:
            screenSize = 'medium'

        #TODO template should be defined from resolution, since that will give the optimal aspect ratio!
        resolution = {}
        resolution['small']  = (800,566)   # legacy screen
        resolution['medium'] = (1024,724)  # normal (old) projector or screen
        resolution['large']  = (1550,1080) # modern projector or screen (~1080p)
        resolution['huge']   = (3055,2160) # high-resolution projektor or screen (~4k)

        if 'screenResolution' in hermodConfiguration.keys():
            screenResolution = map(int,hermodConfiguration['screenResolution'].split('x'))
            resolution['userdefined'] = screenResolution
            screenSize = 'userdefined'

        # Loop over pages
        for page in pages:
            logger.debug('Working with page.')

#            startEpoch = globalStartEpoch

            # TODO: Check filename is correct and do not break security!
            page['filename'] = page['headline'].replace(' ','_') + filetimestampPattern + '.png'

            if 'marginLeft' in page.keys():
                marginLeft = float(page['marginLeft'])
            else:
                marginLeft = 2.0

            if 'marginRight' in page.keys():
                marginRight = float(page['marginRight'])
            else:
                marginRight = 2.0

            if 'marginTop' in page.keys():
                marginTop = float(page['marginTop'])
            else:
                marginTop = 5.0

            if 'marginBottom' in page.keys():
                marginBottom = float(page['marginBottom'])
            else:
                marginBottom = 2.0

#            if 'width' in page.keys():
#                imageWidth = page['width']
#            else:
#                imageWidth = 23

#            if 'height' in page.keys():
#                imageHeight = page['height']
#            else:
#                imageHeight = 12

            if 'ppt_template' in page.keys():
                slideLayout = page['ppt_template']
            else:
                slideLayout = 0

            if 'forecast' in page.keys():
                forecastEnd = page['forecast']
            else:
                forecastEnd = 0

            if 'forecastEnd' in page.keys():
                forecastEnd = page['forecastEnd']
            else:
                forecastEnd = 0

            if 'forecastStart' in page.keys():
                forecastStart = page['forecastStart']
            else:
                forecastStart = 0

            if 'timestep' in page.keys():
                timestep = page['timestep']
            else:
                timestep = 0

            if 'imageScale' in page.keys():
                imageScale = page['imageScale']
            else:
                imageScale = 1.0

            if 'imagePosition' in page.keys():
                imagePosition = map(int,page['imagePosition'].split(','))
            else:
                imagePosition = 1, 1 # Top left corner

            if 'imageSize' in page.keys():
                imageSize = map(int,page['imageSize'].split('x'))
            else:
                imageSize = 6, 4 # Full slide

            if 'imageBorder' in page.keys():
                imageBorder = float(page['imageBorder'])
            else:
                imageBorder = 0.5

            if 'newSlide' in page.keys():
                newSlide = page['newSlide']
            else:
                newSlide = True

            if 'showTitle' in page.keys():
                showTitle = page['showTitle']
            else:
                showTitle = True

            if 'delay' in page.keys():
                delay = page['delay'] * 1000 # convert to microseconds
            else:
                delay = globalDelay

            # Define area for image(s)
            pageBoxWidth    = pageWidth  - marginRight - marginLeft
            pageBoxHeight   = pageHeight - marginBottom - marginTop
            cmToPx          = 0.5 * (resolution[screenSize][0]/pageWidth + resolution[screenSize][1]/pageHeight)
            pxToCm          = 1.0 / cmToPx

            # Define the maximum image resolution for the desired imageSize (will be used when requesting image from url)
            imageResolution = ( int(cmToPx*((imageSize[0]/float(pageLayout[0]))*pageBoxWidth -2*imageBorder)*imageScale) ,
                                int(cmToPx*((imageSize[1]/float(pageLayout[1]))*pageBoxHeight-2*imageBorder)*imageScale) )

            imageWidth  = pxToCm * imageResolution[0] / imageScale
            imageHeight = pxToCm * imageResolution[1] / imageScale
            aspectRatio = float(imageWidth)/float(imageHeight)
            imageAspectRatio = float(imageResolution[0])/float(imageResolution[1])

            # Centering and updating aspect ratio
            imageCenter = ( marginLeft + pageBoxWidth  * (imagePosition[0] - 1 + 0.5*imageSize[0])/float(pageLayout[0]),
                            marginTop  + pageBoxHeight * (imagePosition[1] - 1 + 0.5*imageSize[1])/float(pageLayout[1]) )

            imageLeft = imageCenter[0] - 0.5 * imageWidth
            imageTop  = imageCenter[1] - 0.5 * imageHeight

            # Calculation the individual timesteps for the forecast
            # We always +1 to get the correct result
            # Change: If timestep=0 then just generate forecast time
            if timestep == 0:
#                forecast = 0
                timestep = 1000
            for calculationEpoch in xrange(forecastStart,forecastEnd+1,timestep):
                currentEpoch = startEpoch + datetime.timedelta(hours=calculationEpoch)
                logger.debug('Working with epoch : ' + currentEpoch.strftime('%Y-%m-%dT%H:00:00'))
                currentFilename = formatFilename( page['filename'], currentEpoch )
                currentFilename = get_storage_folder( PowerHermodScratchPath, currentFilename, utcNow)

                # Generate a title on the image
                if showTitle:
                    addedTitleEncoded = page['headline'] + headlinePattern
                else:
                    addedTitleEncoded = ''

                wmsUrl = formatUrl( page['url'], currentEpoch) + '&title=' + urllib.quote( formatUrl( addedTitleEncoded, currentEpoch )  )

                # Timestamps
                tid = formatUrl( '&hidecontrols=ALL&showcontrols=legend&datetime=%Y-%m-%dT%H:00:00', currentEpoch)
                print tid

                url = page['url'] + tid + '&title=' + urllib.quote( formatUrl( addedTitleEncoded, currentEpoch )  )

                # Is this a 'normal' image or a leaflet/openlayers service?
                if 'image' in page.keys():
                    if page['image'] == True:
                        # Get the image, using urllib
                        url = page['url']
                        print url, currentFilename
                        try:
                            urllib.urlretrieve(url, currentFilename)
                        except:
                            currentFilename = ''
                    else:
                        generateImageFromWMS(url, currentFilename,vdisplay.vdisplay_num, imageResolution=imageResolution, delay=delay)
                else:
                    generateImageFromWMS(url, currentFilename,vdisplay.vdisplay_num, imageResolution=imageResolution, delay=delay)

                # Insert the image (generated or fetched on the slide)
                img_path = currentFilename

                # compress images with pngquant
                print compressImages
                if compressImages:
                    subprocess.call('pngquant --force --ext .png ' + img_path, shell=True)

                current_slidelayout = prs.slide_layouts[slideLayout]

                # Only add new slide if we are out of image positions (num_imgs)
                if newSlide:
                    slide = prs.slides.add_slide(current_slidelayout)

#                titleBox = slide.shapes.add_textbox(Cm(5), Cm(2), Cm(10), Cm(10))
#                titleFrame = titleBox.text_frame
                    title = slide.shapes.title
                    slideTitle = page['headline']# + headlinePattern
                    title.text = formatTitle( slideTitle, currentEpoch )
#                titleFrame.text = formatTitle( slideTitle, currentEpoch )
#                    title.font.size = Pt(40)

                left   = Cm(imageLeft)
                top    = Cm(imageTop)
                width  = Cm(imageWidth)
                height = Cm(imageHeight)

                # Check if images are created
                if os.path.isfile(img_path):
                        # Does it have a size that looks like it is a image
                        #TODO: Use PIL to verify that it is a image
                        if os.path.getsize(img_path) > 1000:
                            pic = slide.shapes.add_picture(img_path,left,top,width,height)
                            line = pic.line
                            line.color.rgb = RGBColor(128,128,128)
                            line.width = Cm(0.02)
                        else:
                            width = height = Cm(10)
                            top = Cm(10)
                            left =  Cm(5)
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.textframe
                            tf.text = hermodConfiguration['missingimagetext']


        outputFilename = formatFilename(outputFilename, utcNow)
        powerpointFilename = get_storage_folder( PowerHermodScratchPath, outputFilename, utcNow)
        prs.save( powerpointFilename )
        powerpointFilename = get_storage_folder( PowerHermodArchivePath, outputFilename, utcNow)
        prs.save( powerpointFilename )
        vdisplay.stop()


        # Create message container

        # Create title for the email
        titlePlain = pptTitle + headlinePattern
        titlePlain = formatTitle( titlePlain.encode('utf-8'), utcNow )
        fromEmail = 'info@fcoo.dk'

        msg = MIMEMultipart()
        msg['Subject'] = titlePlain.replace('_',' ') + " prognose fra fcoo.dk"
        msg['From'] = fromEmail

        # Create the body of the message (a plain-text).
        # TODO remove titlePlain hack!
        text = "Prognosen for " + titlePlain + " fra FCOO vedlagt som powerpoint dokument. Kontakt info@fcoo.dk for yderlig information.\n\n"
        text = text + "Forecast for " + titlePlain + " from FCOO as a powerpoint attachment. Contact information : info@fcoo.dk \n\n"
        text = text + emailtext

        # Record the MIME types of both parts - text/plain and text/html.
        textPart = MIMEText(text.encode('utf-8'), 'plain')
        powerpointPart = MIMEBase('application', 'vnd.ms-powerpoint')

        powerpointPart.set_payload( open( powerpointFilename ,"rb").read() )
        Encoders.encode_base64( powerpointPart )
        powerpointPart.add_header('Content-Disposition', 'attachment; filename="%s"' % os.path.basename(powerpointFilename))

        # Attach parts into message container.
        # According to RFC 2046, the last part of a multipart message, in this case
        # the HTML message, is best and preferred.
        msg.attach(textPart)
        msg.attach(powerpointPart)

        disableSMTP = False
        smtpserver = '127.0.0.1'
        
        # Handle recipents
        msgTo = readAndVerifyEmailRecipient( hermodConfiguration['recipient'] )
        logger.info("Recipient : " + str(msgTo))
                                        
        if disableSMTP <> True:
            # Send the message via a SMTP server
            smtpConnection = smtplib.SMTP(smtpserver)

            # sendmail function takes 3 arguments: sender's address, recipient's address
            # and message to send - here it is sent as one string.
            smtpConnection.sendmail(fromEmail, msgTo, msg.as_string())
            smtpConnection.quit()


    # No epoch match the current time!
    else:
        logger.info('No match on current time')

class CLIError(Exception):
    '''Generic exception to raise and log different fatal errors.'''
    def __init__(self, msg):
        super(CLIError).__init__(type(self))
        self.msg = "E: %s" % msg
    def __str__(self):
        return self.msg
    def __unicode__(self):
        return self.msg

def main(argv=None): # IGNORE:C0111
    '''Command line options.'''

    if argv is None:
        argv = sys.argv
    else:
        sys.argv.extend(argv)

    program_name = os.path.basename(sys.argv[0])
    program_version = "v%s" % __version__
    program_build_date = str(__updated__)
    program_version_message = '%%(prog)s %s (%s)' % (program_version, program_build_date)
    program_shortdesc = __import__('__main__').__doc__.split("\n")[1]
    program_license = '''%s

  Created by user_name on %s.
  Copyright 2014 FCOO. All rights reserved.

  Licensed under the Apache License 2.0
  http://www.apache.org/licenses/LICENSE-2.0

  Distributed on an "AS IS" basis without warranties
  or conditions of any kind, either express or implied.

USAGE
''' % (program_shortdesc, str(__date__))

    try:
        # Setup argument parser
        parser = ArgumentParser(description=program_license, formatter_class=RawDescriptionHelpFormatter)
        parser.add_argument("-v", "--verbose", dest="verbose", action="count", help="set verbosity level [default: %(default)s]")
        parser.add_argument("-i", "--include", dest="include", help="only include paths matching this regex pattern. Note: exclude is given preference over include. [default: %(default)s]", metavar="RE" )
        parser.add_argument('-V', '--version', action='version', version=program_version_message)
        parser.add_argument(dest="paths", help="paths to folder(s) with source file(s) [default: %(default)s]", metavar="path", nargs='+')

        # Process arguments
        args = parser.parse_args()

        paths = args.paths
        verbose = args.verbose
        recurse = args.recurse
        inpat = args.include
        expat = args.exclude

        if verbose > 0:
            print("Verbose mode on")
            if recurse:
                print("Recursive mode on")
            else:
                print("Recursive mode off")

        if inpat and expat and inpat == expat:
            raise CLIError("include and exclude pattern are equal! Nothing will be processed.")

        for inpath in paths:
            ### do something with inpath ###
            print(inpath)
        return 0
    except KeyboardInterrupt:
        ### handle keyboard interrupt ###
        return 0
    except Exception, e:
        if DEBUG or TESTRUN:
            raise(e)
        indent = len(program_name) * " "
        sys.stderr.write(program_name + ": " + repr(e) + "\n")
        sys.stderr.write(indent + "  for help use --help")
        return 2

if __name__ == "__main__":

    # Setup logger
    logfmt = "%(asctime)-15s %(levelname)s %(message)s"
    loglevel = logging.INFO
    logger = logging.getLogger()
    formatter = logging.Formatter(logfmt)
    ch = logging.StreamHandler(sys.stdout)
    ch.setLevel(loglevel)
    ch.setFormatter(formatter)
    logger.addHandler(ch)
    logger.setLevel(loglevel)

    # Start loading the basic Grane konfiguration
    graneYamlStream = file('Grane.yaml','r')
    graneConfiguration = yaml.load( graneYamlStream )

    powerhermodConfigurationFiles = graneConfiguration['powerhermodConfigurations']

    PowerHermodConfigurationFileStorage = os.path.join(os.path.normpath(powerhermodConfigurationFiles), "*.yaml")

    for fname in glob.glob(PowerHermodConfigurationFileStorage):
        print fname
        stream = file(fname, 'r')
        hermodConfiguration = yaml.load(stream)
        processHermodSubscription( hermodConfiguration ) 
